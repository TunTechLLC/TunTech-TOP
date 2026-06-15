"""Tests for document_processor utility functions.
These tests do not call Claude — they test local processing logic only.
"""
import io
import os
import tempfile
import pytest


def test_markdown_fence_stripping():
    """Verify JSON wrapped in code fences is cleaned correctly.
    Tests all fence variants Claude might produce despite being told not to."""
    from api.services.document_processor import strip_json_fences

    raw_json = '[{"signal_name": "test", "domain": "Delivery Operations"}]'

    # No fences — should pass through unchanged
    assert strip_json_fences(raw_json) == raw_json

    # ```json fence
    fenced_json = f'```json\n{raw_json}\n```'
    assert strip_json_fences(fenced_json) == raw_json

    # plain ``` fence
    fenced_plain = f'```\n{raw_json}\n```'
    assert strip_json_fences(fenced_plain) == raw_json

    # Extra whitespace
    fenced_whitespace = f'  ```json\n{raw_json}\n```  '
    assert strip_json_fences(fenced_whitespace) == raw_json


def test_file_type_detection():
    """Verify get_file_type correctly identifies file types from filename."""
    from api.services.document_processor import get_file_type

    assert get_file_type('E001_interview_CEO_Okafor.txt')        == 'interview'
    assert get_file_type('E001_financial_FY2025_PL.txt')         == 'financial'
    assert get_file_type('E001_sow_ProjectAlpha.txt')            == 'sow'
    assert get_file_type('E001_portfolio_March.txt')             == 'portfolio'
    assert get_file_type('E001_status_ProjectAlpha_March.txt')   == 'status'
    assert get_file_type('E001_resource_utilization_Q1.txt')     == 'resource'
    assert get_file_type('E001_delivery_risk_register.txt')      == 'delivery'
    assert get_file_type('E001_other_misc.txt')                  == 'other'
    assert get_file_type('E001_unknown_type_file.txt')           == 'other'
    # Type with uppercase — should normalize to lowercase
    assert get_file_type('E001_INTERVIEW_CEO.txt')               == 'interview'


def test_file_type_detection_new_convention():
    """Verify get_file_type works with Interview_/Doc_ convention and any extension."""
    from api.services.document_processor import get_file_type

    # Interview_ prefix — always 'interview' regardless of extension
    assert get_file_type('Interview_CEO.txt')              == 'interview'
    assert get_file_type('Interview_CEO.docx')             == 'interview'
    assert get_file_type('Interview_DirectorDelivery.pdf') == 'interview'

    # Doc_ prefix — stem determines type, extension irrelevant
    assert get_file_type('Doc_Financial.xlsx')             == 'financial'
    assert get_file_type('Doc_Financial_Q1.xlsx')          == 'financial'
    assert get_file_type('Doc_SOW.pdf')                    == 'sow'
    assert get_file_type('Doc_SOW_ProjectAlpha.pdf')       == 'sow'
    assert get_file_type('Doc_Portfolio.pptx')             == 'portfolio'
    assert get_file_type('Doc_StatusReport.pdf')           == 'status'
    assert get_file_type('Doc_Resource.xlsx')              == 'resource'
    assert get_file_type('Doc_Delivery.docx')              == 'delivery'
    assert get_file_type('Doc_Other.pdf')                  == 'other'


def test_supported_extensions():
    """Verify SUPPORTED_EXTENSIONS contains exactly the expected formats."""
    from api.services.document_processor import SUPPORTED_EXTENSIONS

    assert SUPPORTED_EXTENSIONS == {'.txt', '.docx', '.xlsx', '.pdf', '.pptx'}


def test_valid_domains_in_utils():
    """Verify domains.py contains exactly the expected 10 domains."""
    from api.utils.domains import VALID_DOMAINS

    expected = {
        'Sales & Pipeline',
        'Sales-to-Delivery Transition',
        'Delivery Operations',
        'Resource Management',
        'Project Governance / PMO',
        'Consulting Economics',
        'Customer Experience',
        'AI Readiness',
        'Human Resources',
        'Finance and Commercial',
    }
    assert VALID_DOMAINS == expected, (
        f"Domain mismatch: {VALID_DOMAINS.symmetric_difference(expected)}"
    )


def test_valid_confidences_in_utils():
    """Verify domains.py contains the correct confidence values."""
    from api.utils.domains import VALID_CONFIDENCES

    assert 'High' in VALID_CONFIDENCES
    assert 'Medium' in VALID_CONFIDENCES
    assert 'Hypothesis' in VALID_CONFIDENCES
    assert len(VALID_CONFIDENCES) == 3


# ---------------------------------------------------------------------------
# extract_text_from_file — unit tests using programmatically-created files
# No engagement, no Claude, no DB required.
# ---------------------------------------------------------------------------

def test_extract_txt_returns_content():
    from api.services.document_processor import extract_text_from_file
    with tempfile.NamedTemporaryFile(mode='w', suffix='.txt', delete=False,
                                    encoding='utf-8') as f:
        f.write('This is a test transcript.\nSecond line.')
        path = f.name
    try:
        result = extract_text_from_file(path, 'Interview_CEO.txt')
        assert 'test transcript' in result
        assert 'Second line' in result
    finally:
        os.unlink(path)


def test_extract_txt_empty_raises():
    from api.services.document_processor import extract_text_from_file
    with tempfile.NamedTemporaryFile(mode='w', suffix='.txt', delete=False,
                                    encoding='utf-8') as f:
        f.write('   \n  ')
        path = f.name
    try:
        with pytest.raises(ValueError, match='empty'):
            extract_text_from_file(path, 'Interview_CEO.txt')
    finally:
        os.unlink(path)


def test_extract_unsupported_extension_raises():
    from api.services.document_processor import extract_text_from_file
    with tempfile.NamedTemporaryFile(suffix='.csv', delete=False) as f:
        f.write(b'col1,col2\nval1,val2\n')
        path = f.name
    try:
        with pytest.raises(ValueError, match='unsupported'):
            extract_text_from_file(path, 'data.csv')
    finally:
        os.unlink(path)


def test_extract_docx_paragraphs_and_tables():
    """Word document: verify both paragraph text and table cell text are extracted."""
    pytest.importorskip('docx', reason='python-docx not installed')
    from docx import Document
    from api.services.document_processor import extract_text_from_file

    doc = Document()
    doc.add_paragraph('CEO stated delivery margins are under pressure.')
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = 'Project'
    table.cell(0, 1).text = 'Status'
    table.cell(1, 0).text = 'Alpha'
    table.cell(1, 1).text = 'At risk'

    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as f:
        path = f.name
    doc.save(path)

    try:
        result = extract_text_from_file(path, 'Interview_CEO.docx')
        assert 'delivery margins' in result
        assert 'At risk' in result
        assert 'Project' in result
    finally:
        os.unlink(path)


def test_extract_docx_empty_raises():
    pytest.importorskip('docx', reason='python-docx not installed')
    from docx import Document
    from api.services.document_processor import extract_text_from_file

    doc = Document()
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as f:
        path = f.name
    doc.save(path)

    try:
        with pytest.raises(ValueError, match='no extractable text'):
            extract_text_from_file(path, 'Doc_SOW.docx')
    finally:
        os.unlink(path)


def test_extract_xlsx_sheets_and_rows():
    """Excel: verify sheet headers and row data are extracted."""
    pytest.importorskip('openpyxl', reason='openpyxl not installed')
    import openpyxl
    from api.services.document_processor import extract_text_from_file

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = 'P&L'
    ws.append(['Metric', 'Value'])
    ws.append(['Gross Margin', '42%'])
    ws.append(['Revenue', '2400000'])

    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as f:
        path = f.name
    wb.save(path)

    try:
        result = extract_text_from_file(path, 'Doc_Financial.xlsx')
        assert '[Sheet: P&L]' in result
        assert 'Gross Margin' in result
        assert '42%' in result
    finally:
        os.unlink(path)


def test_extract_xlsx_empty_raises():
    pytest.importorskip('openpyxl', reason='openpyxl not installed')
    import openpyxl
    from api.services.document_processor import extract_text_from_file

    wb = openpyxl.Workbook()
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as f:
        path = f.name
    wb.save(path)

    try:
        with pytest.raises(ValueError, match='no extractable data'):
            extract_text_from_file(path, 'Doc_Financial.xlsx')
    finally:
        os.unlink(path)


def test_extract_pptx_slides_and_notes():
    """PowerPoint: verify slide text and speaker notes are extracted."""
    pytest.importorskip('pptx', reason='python-pptx not installed')
    from pptx import Presentation
    from pptx.util import Inches
    from api.services.document_processor import extract_text_from_file

    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = 'Delivery Overview'
    slide.placeholders[1].text = 'Three projects are currently at risk.'

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = 'CEO confirmed this in the follow-up call.'

    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
        path = f.name
    prs.save(path)

    try:
        result = extract_text_from_file(path, 'Doc_StatusReport.pptx')
        assert '[Slide 1]' in result
        assert 'Delivery Overview' in result
        assert 'at risk' in result
        assert 'CEO confirmed' in result
    finally:
        os.unlink(path)


def test_extract_pptx_empty_raises():
    pytest.importorskip('pptx', reason='python-pptx not installed')
    from pptx import Presentation
    from api.services.document_processor import extract_text_from_file

    prs = Presentation()
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
        path = f.name
    prs.save(path)

    try:
        with pytest.raises(ValueError, match='no extractable text'):
            extract_text_from_file(path, 'Doc_StatusReport.pptx')
    finally:
        os.unlink(path)


def test_archive_candidate_files_new_convention():
    """archive_candidate_files() must archive Interview_/Doc_ named candidate files,
    not just legacy engagement_id-prefixed ones."""
    import json
    from api.services.document_processor import archive_candidate_files

    with tempfile.TemporaryDirectory() as candidates_folder:
        # Merged file
        merged = os.path.join(candidates_folder, 'E005_merged_candidates.json')
        # Individual files — new naming convention (no engagement_id prefix)
        interview_file = os.path.join(candidates_folder, 'Interview_CEO_candidates.json')
        doc_file = os.path.join(candidates_folder, 'Doc_Financial_candidates.json')

        for path in [merged, interview_file, doc_file]:
            with open(path, 'w') as f:
                json.dump({}, f)

        archive_candidate_files('E005', candidates_folder, merged)

        processed_dir = os.path.join(candidates_folder, 'processed')
        archived = os.listdir(processed_dir)

        assert 'E005_merged_candidates.json' in archived
        assert 'Interview_CEO_candidates.json' in archived
        assert 'Doc_Financial_candidates.json' in archived
        # All three moved — none left in candidates_folder root
        assert not os.path.exists(merged)
        assert not os.path.exists(interview_file)
        assert not os.path.exists(doc_file)


# --- Defect A: JSON robustness (pure parser raises, retry recovers, never silent-empty) ---

def test_parse_extraction_response_valid_dict():
    """New {found, not_observed} shape parses to (candidates, not_observed)."""
    from api.services.document_processor import _parse_extraction_response
    raw = '{"found": [{"signal_name": "x"}], "not_observed": ["SL-1"]}'
    candidates, not_observed = _parse_extraction_response(raw, "f.txt")
    assert len(candidates) == 1 and not_observed == ["SL-1"]


def test_parse_extraction_response_legacy_list():
    """Legacy bare-array shape parses to (list, [])."""
    from api.services.document_processor import _parse_extraction_response
    candidates, not_observed = _parse_extraction_response('[{"signal_name": "y"}]', "f.txt")
    assert len(candidates) == 1 and not_observed == []


def test_parse_extraction_response_raises_on_malformed_not_silent():
    """Core Defect A guarantee: malformed JSON RAISES — never returns empty silently."""
    import json
    from api.services.document_processor import _parse_extraction_response
    with pytest.raises(json.JSONDecodeError):
        _parse_extraction_response('{"found": [ {"a":1} {"b":2} ]}', "f.txt")  # missing comma
    with pytest.raises(json.JSONDecodeError):
        _parse_extraction_response('not json at all', "f.txt")


def test_parse_extraction_response_raises_on_unexpected_shape():
    """A valid-JSON-but-wrong-shape response (e.g. a bare number) raises ValueError."""
    from api.services.document_processor import _parse_extraction_response
    with pytest.raises(ValueError):
        _parse_extraction_response('42', "f.txt")


def test_extract_signals_with_retry_recovers_on_second_attempt(monkeypatch):
    """A transient malformed response is recovered by the single retry (regeneration)."""
    import asyncio
    from api.services import document_processor as dp
    import api.services.claude as claude_mod

    calls = {"n": 0}
    async def fake(content, library_block):
        calls["n"] += 1
        if calls["n"] == 1:
            return '{"found": [ {"a":1} {"b":2} ]}'   # malformed first
        return '{"found": [{"signal_name": "ok"}], "not_observed": []}'  # valid on retry
    monkeypatch.setattr(claude_mod, "extract_signals_from_transcript", fake)

    candidates, not_observed = asyncio.run(
        dp._extract_signals_with_retry("interview", "text", None, "", "f.txt")
    )
    assert calls["n"] == 2          # it retried
    assert len(candidates) == 1     # and recovered


def test_extract_signals_with_retry_raises_after_all_attempts(monkeypatch):
    """When every attempt is unparseable, it RAISES (surfaces) — never silent-empty.
    The caller records the file failed and leaves it unprocessed for re-attempt."""
    import asyncio
    from api.services import document_processor as dp
    import api.services.claude as claude_mod

    async def always_bad(content, library_block):
        return "not valid json"
    monkeypatch.setattr(claude_mod, "extract_signals_from_transcript", always_bad)

    with pytest.raises(ValueError):
        asyncio.run(dp._extract_signals_with_retry("interview", "text", None, "", "f.txt"))


# --- Defect B: semantic-dedup merge, corroboration, fallback, within-domain invariant ---

def test_merge_cluster_single_unchanged():
    from api.services.document_processor import _merge_cluster
    out = _merge_cluster([{"signal_name": "x", "signal_confidence": "Medium", "source_file": "f"}])
    assert out["signal_name"] == "x" and out["signal_confidence"] == "Medium"
    assert out["_corroboration"] == 1          # singleton is corroborated by its one source


def test_merge_cluster_corroboration_upgrade_two_sources():
    """A signal corroborated by >=2 distinct sources is upgraded one level + noted."""
    from api.services.document_processor import _merge_cluster
    m = _merge_cluster([
        {"signal_name": "a", "signal_confidence": "Hypothesis", "source_file": "f1.txt", "notes": "n"},
        {"signal_name": "b", "signal_confidence": "Medium", "source_file": "f2.txt", "notes": "longer note"},
    ])
    assert m["signal_confidence"] == "High"            # max(Medium,Hyp)=Medium, +1 -> High
    assert "Corroborated across 2 sources" in m["notes"]


def test_merge_cluster_no_upgrade_single_source():
    """Multiple restatements from the SAME file are not corroboration — no upgrade."""
    from api.services.document_processor import _merge_cluster
    m = _merge_cluster([
        {"signal_name": "a", "signal_confidence": "Hypothesis", "source_file": "f1.txt"},
        {"signal_name": "b", "signal_confidence": "Medium", "source_file": "f1.txt"},
    ])
    assert m["signal_confidence"] == "Medium"


def test_consolidate_fallback_preserves_all_and_flags(monkeypatch):
    """If semantic clustering fails, fall back to exact-only — nothing lost, degraded=True."""
    import asyncio
    from api.services import document_processor as dp
    import api.services.claude as claude_mod

    async def boom(items):
        raise RuntimeError("simulated clustering failure")
    monkeypatch.setattr(claude_mod, "cluster_duplicate_signals", boom)

    cands = [
        {"domain": "D", "signal_name": "S1", "signal_confidence": "High", "source_file": "a"},
        {"domain": "D", "signal_name": "S1", "signal_confidence": "Medium", "source_file": "b"},
        {"domain": "D", "signal_name": "S2", "signal_confidence": "High", "source_file": "a"},
    ]
    out, removed, degraded = asyncio.run(dp._consolidate_candidates(cands))
    assert degraded is True
    assert len(out) == 2 and removed == 1          # S1 (exact-merged) + S2; nothing lost


def test_consolidate_enforces_within_domain(monkeypatch):
    """Even if the model groups across domains, code splits the cluster by domain."""
    import asyncio
    from api.services import document_processor as dp
    import api.services.claude as claude_mod

    async def cross_domain(items):
        return [[0, 1]]                              # model wrongly groups across domains
    monkeypatch.setattr(claude_mod, "cluster_duplicate_signals", cross_domain)

    cands = [
        {"domain": "Sales", "signal_name": "Concentration", "signal_confidence": "High", "source_file": "a"},
        {"domain": "Finance", "signal_name": "Concentration", "signal_confidence": "High", "source_file": "b"},
    ]
    out, _, _ = asyncio.run(dp._consolidate_candidates(cands))
    assert len(out) == 2                            # not merged across domains
    assert {c["domain"] for c in out} == {"Sales", "Finance"}
