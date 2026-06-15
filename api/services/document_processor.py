import os
import json
import logging
import hashlib
import shutil
from pathlib import Path
from datetime import date

from api.db.repositories.processed_files import ProcessedFilesRepository
from api.utils.domains import DEFAULT_DOMAIN, VALID_DOMAINS, VALID_CONFIDENCES, VALID_SIGNAL_VALENCES

logger = logging.getLogger(__name__)

VALID_FILE_TYPES = {'interview', 'financial', 'portfolio', 'sow', 'status', 'resource', 'delivery', 'other'}
SUPPORTED_EXTENSIONS = {'.txt', '.docx', '.xlsx', '.pdf', '.pptx'}

# Domain list injected into every extraction prompt — sourced from VALID_DOMAINS, not hardcoded
_DOMAIN_LIST = ', '.join(f'"{d}"' for d in sorted(VALID_DOMAINS))

# Which domains to check per document file type (Tier 1 library slice)
DOMAIN_FILTER_MAP = {
    'financial': ['Consulting Economics', 'Finance and Commercial'],
    'resource':  ['Resource Management', 'Consulting Economics'],
    'delivery':  ['Delivery Operations', 'Project Governance / PMO',
                  'Sales-to-Delivery Transition', 'Resource Management'],
    'portfolio': ['Delivery Operations', 'Resource Management',
                  'Project Governance / PMO', 'Sales-to-Delivery Transition'],
    'sow':       ['Sales-to-Delivery Transition', 'Finance and Commercial'],
    'status':    ['Delivery Operations', 'Project Governance / PMO'],
    'other':     sorted(VALID_DOMAINS),
    'interview': sorted(VALID_DOMAINS),
}

# Shared library instruction and output format appended to every extraction prompt
_LIBRARY_INSTRUCTION = """SIGNAL LIBRARY:
The user message includes a SIGNAL LIBRARY block listing Tier 1 signals to check against this source.
- For each listed signal where you find evidence: include it in "found" with all required fields plus "library_signal_id": "<SL-XX>"
- For each listed signal you actively checked but found no evidence for: add its signal_id to "not_observed"
- You may include freely-extracted signals not in the library in "found" — omit library_signal_id for these
- Report not_observed ONLY for signals that appear in the SIGNAL LIBRARY block above

"""

_OUTPUT_FORMAT_OBJECT = """CRITICAL OUTPUT FORMAT:
Your response must be a JSON object beginning with { and ending with }
Do not include any text, explanation, or markdown before or after the JSON object
Do not use code fences or backticks of any kind
Your response must follow this structure exactly:
{"found": [...signal objects...], "not_observed": ["SL-XX", ...]}
The "found" array contains extracted signal objects following all rules above.
The "not_observed" array contains signal_ids from the SIGNAL LIBRARY block that you actively checked and could not find evidence for."""

# Shared valence field spec — positive evidence enters the pipeline at extraction.
# Inserted into every document extraction prompt's field list so a healthy metric is
# captured as a Strength, not only problems as Risks.
_VALENCE_FIELD = """- valence: string — the signal's directional meaning, exactly one of:
    - "Strength": evidence of operational health — a metric, capability, or behavior working well and worth preserving (e.g. margin above benchmark, strong retention, a disciplined process)
    - "Risk": evidence of dysfunction, weakness, or a negative indicator (e.g. margin erosion, missed milestones, unmitigated risks)
    - "Neutral": a contextual fact that is neither clearly positive nor negative
  When the evidence is genuinely ambiguous, use "Neutral". Do not label something a Strength without specific evidence — an unsupported strength is as invalid as an unsupported problem."""

FINANCIAL_EXTRACTION_PROMPT = f"""You are analyzing a financial document from a consulting firm diagnostic engagement.

Extract signals that are directly supported by data in the document. Focus on economics, margin, utilization, revenue, pricing, and financial health indicators.

Financial documents may contain tables of numbers with minimal narrative context.
When you see numerical data, infer the signal from the numbers themselves.
A gross margin percentage in a row labeled 'Gross Margin' is directly readable —
use the row label as the signal name and apply the confidence rules below.

Each item in "found" must have exactly these fields:
- signal_name: string
- domain: string — must be exactly one of: {_DOMAIN_LIST}
- observed_value: string
- normalized_band: string
- evidence_quality: string — COMPLETE THIS BEFORE assigning signal_confidence. Scan the document for any issues with this signal's evidence: contradictory figures elsewhere in the document for this metric; whether the figure is inferred by you rather than explicitly stated; missing labels or context that make interpretation ambiguous. Write "None" only if the evidence is explicitly stated, clearly labeled, and unambiguous.
- signal_confidence: string — derived from evidence_quality:
    - Hypothesis: if evidence_quality contains anything other than "None"
    - Medium: if evidence_quality is "None" AND the signal is qualitative with no specific number
    - High: if evidence_quality is "None" AND a specific number is explicitly stated with a clear label
- source: string — always "Document"
- economic_relevance: string
- notes: string — include the specific data point from the document that supports this signal
{_VALENCE_FIELD}
- library_signal_id: string — ONLY include when this signal matches an entry from the SIGNAL LIBRARY block. Use the exact signal_id (e.g. "SL-38"). Omit for freely-extracted signals.

Only extract signals with direct document evidence.
Extract no more than 10 found signals. If you identify more, keep only the 10 most operationally significant.

""" + _LIBRARY_INSTRUCTION + _OUTPUT_FORMAT_OBJECT

PORTFOLIO_EXTRACTION_PROMPT = f"""You are analyzing a project portfolio or status report from a consulting firm diagnostic engagement.

Extract signals related to delivery performance, project health, governance, and resource utilization.

Look specifically for:
- RAG status indicators (Red/Amber/Green) on individual projects
- Schedule variance — planned vs actual dates
- Budget variance — planned vs actual cost
- Resource utilization percentages
- Issues and risks listed in the report
- Projects described as 'at risk', 'delayed', or 'escalated'

Each item in "found" must have exactly these fields:
- signal_name: string
- domain: string — must be exactly one of: {_DOMAIN_LIST}
- observed_value: string
- normalized_band: string
- evidence_quality: string — COMPLETE THIS BEFORE assigning signal_confidence. Scan the document for any issues with this signal's evidence: contradictory figures elsewhere in the document for this metric; whether the figure is inferred by you rather than explicitly stated; missing labels or context that make interpretation ambiguous. Write "None" only if the evidence is explicitly stated, clearly labeled, and unambiguous.
- signal_confidence: string — derived from evidence_quality:
    - Hypothesis: if evidence_quality contains anything other than "None"
    - Medium: if evidence_quality is "None" AND the signal is qualitative with no specific number
    - High: if evidence_quality is "None" AND a specific number is explicitly stated with a clear label
- source: string — always "Document"
- economic_relevance: string
- notes: string — include the specific data point from the document
{_VALENCE_FIELD}
- library_signal_id: string — ONLY include when this signal matches an entry from the SIGNAL LIBRARY block. Use the exact signal_id (e.g. "SL-17"). Omit for freely-extracted signals.

Extract no more than 10 found signals. If you identify more, keep only the 10 most operationally significant.

""" + _LIBRARY_INSTRUCTION + _OUTPUT_FORMAT_OBJECT

SOW_EXTRACTION_PROMPT = f"""You are analyzing a Statement of Work or contract from a consulting firm diagnostic engagement.

Extract signals related to SOW quality, scope definition, acceptance criteria, assumptions, and sales-to-delivery transition indicators.

Focus particularly on:
- Presence or absence of measurable acceptance criteria
- Clarity of deliverable definitions
- Inclusion of an assumptions section
- Change control provisions
- Whether delivery team involvement is evident from document language
- Success criteria that are measurable vs vague

Each item in "found" must have exactly these fields:
- signal_name: string
- domain: string — must be exactly one of: {_DOMAIN_LIST}
- observed_value: string
- normalized_band: string
- evidence_quality: string — COMPLETE THIS BEFORE assigning signal_confidence. Note any of the following: the signal is an inference from document tone rather than an explicit statement; contradictory language about the same element appears elsewhere in the document; the evidence is a single ambiguous phrase requiring interpretation. Write "None" only if the signal is unambiguously observable without interpretation.
- signal_confidence: string — derived from evidence_quality:
    - Hypothesis: if evidence_quality contains anything other than "None"
    - Medium: if evidence_quality is "None" AND the signal is qualitative with no specific number
    - High: if evidence_quality is "None" AND the presence or absence of a structural element is directly and unambiguously demonstrable
- source: string — always "Document"
- economic_relevance: string
- notes: string — include the specific language from the SOW that supports this signal
{_VALENCE_FIELD}
- library_signal_id: string — ONLY include when this signal matches an entry from the SIGNAL LIBRARY block. Use the exact signal_id (e.g. "SL-13"). Omit for freely-extracted signals.

Extract no more than 10 found signals. If you identify more, keep only the 10 most operationally significant.

""" + _LIBRARY_INSTRUCTION + _OUTPUT_FORMAT_OBJECT

STATUS_EXTRACTION_PROMPT = f"""You are analyzing a project status report from a consulting firm diagnostic engagement.

Extract signals related to delivery health, project performance, and governance quality.

Focus on:
- Schedule status — on track, delayed, at risk, milestone completions
- Budget status — on budget, over budget, variance amounts
- Issues and risks — are they listed with owners and resolution dates?
- Escalations — what has been escalated and to whom?
- Action items — are they tracked with owners and due dates?
- Overall project health indicators — RAG status, confidence levels

Return only signals that are clearly observable in the document. Do not infer signals
that are not supported by specific text.
Extract no more than 10 found signals. If you identify more, keep only the 10 most operationally significant.

Each item in "found" must have exactly these fields:
- signal_name: string — short descriptive name (e.g., "Multiple projects Red status")
- domain: string — must be exactly one of: {_DOMAIN_LIST}
- observed_value: string — the specific value or observation (e.g., "3 of 7 projects Red")
- normalized_band: string — context for the value (e.g., "Above acceptable threshold")
- evidence_quality: string — COMPLETE THIS BEFORE assigning signal_confidence. Scan the document for any issues with this signal's evidence: contradictory figures elsewhere in the document for this metric; whether the figure is inferred by you rather than explicitly stated; missing labels or context that make interpretation ambiguous. Write "None" only if the evidence is explicitly stated, clearly labeled, and unambiguous.
- signal_confidence: string — derived from evidence_quality:
    - Hypothesis: if evidence_quality contains anything other than "None"
    - Medium: if evidence_quality is "None" AND the signal is qualitative with no specific number
    - High: if evidence_quality is "None" AND a specific number is explicitly stated with a clear label
- source: string — always "Document"
- economic_relevance: string — brief note on economic impact, or empty string
- notes: string — direct quote or specific reference from the document
{_VALENCE_FIELD}
- library_signal_id: string — ONLY include when this signal matches an entry from the SIGNAL LIBRARY block. Use the exact signal_id (e.g. "SL-36"). Omit for freely-extracted signals.

""" + _LIBRARY_INSTRUCTION + _OUTPUT_FORMAT_OBJECT

RESOURCE_EXTRACTION_PROMPT = f"""You are analyzing a resource utilization or staffing report from a consulting firm diagnostic engagement.

Extract signals related to utilization rates, staffing patterns, capacity, and resource management.

Focus on:
- Overall utilization percentage and target (e.g., 71% actual vs 78% target)
- Individual consultant utilization spread — are some overloaded while others are underutilized?
- Bench time — how many consultants are unbilled and for how long?
- Staffing assignment patterns — are projects understaffed or overstaffed?
- Hiring activity — open positions, time-to-fill, attrition
- Capacity vs pipeline alignment — do you have the right skills for upcoming work?

Return only signals that are clearly observable in the document.
Extract no more than 10 found signals. If you identify more, keep only the 10 most operationally significant.

Each item in "found" must have exactly these fields:
- signal_name: string — short descriptive name
- domain: string — must be exactly one of: {_DOMAIN_LIST}
- observed_value: string — the specific value (e.g., "71% billable utilization")
- normalized_band: string — context (e.g., "Below 78% target")
- evidence_quality: string — COMPLETE THIS BEFORE assigning signal_confidence. Scan the document for any issues with this signal's evidence: contradictory figures elsewhere in the document for this metric; whether the figure is inferred by you rather than explicitly stated; missing labels or context that make interpretation ambiguous. Write "None" only if the evidence is explicitly stated, clearly labeled, and unambiguous.
- signal_confidence: string — derived from evidence_quality:
    - Hypothesis: if evidence_quality contains anything other than "None"
    - Medium: if evidence_quality is "None" AND the signal is qualitative with no specific number
    - High: if evidence_quality is "None" AND a specific number is explicitly stated with a clear label
- source: string — always "Document"
- economic_relevance: string — brief note on economic impact, or empty string
- notes: string — direct reference from the document
{_VALENCE_FIELD}
- library_signal_id: string — ONLY include when this signal matches an entry from the SIGNAL LIBRARY block. Use the exact signal_id (e.g. "SL-25"). Omit for freely-extracted signals.

""" + _LIBRARY_INSTRUCTION + _OUTPUT_FORMAT_OBJECT

DELIVERY_DOCUMENT_EXTRACTION_PROMPT = f"""You are analyzing a delivery document from a consulting firm diagnostic engagement.

Delivery documents include risk registers, project retrospectives, portfolio summaries, and proposals.
Extract signals that are directly supported by evidence in the document.

Focus on:

Risk registers:
- Open risk count and severity distribution — how many High/Medium/Low risks are open?
- Mitigation coverage — are risks listed without mitigations or owners?
- Risk ownership gaps — risks with no assigned owner
- Escalated risks — items flagged for management attention

Retrospectives:
- Recurring issues — problems that appear across multiple retrospectives
- Action item completion — were prior retro actions actually resolved?
- Process improvement adoption — were changes from past retros implemented?
- Team health indicators — morale, capacity, or workload signals in retro language

Portfolio summaries:
- Aggregate delivery health — overall on-track / at-risk / red counts across the portfolio
- Capacity vs pipeline alignment — is current staffing sufficient for the active portfolio?
- Revenue recognition patterns — delayed milestones affecting invoicing

Proposals:
- Scope definition quality — are deliverables specific and measurable, or vague?
- Delivery methodology — is how delivery will be executed described, or absent?
- Assumptions coverage — is there an explicit assumptions section?
- Pricing structure — fixed fee, T&M, or hybrid; any margin risk indicators in the structure

Return only signals that are clearly observable in the document.
Extract no more than 10 found signals. If you identify more, keep only the 10 most operationally significant.

Each item in "found" must have exactly these fields:
- signal_name: string — short descriptive name (e.g., "High open risk count with no mitigations")
- domain: string — must be exactly one of: {_DOMAIN_LIST}
- observed_value: string — the specific observation (e.g., "14 of 22 risks unmitigated")
- normalized_band: string — context for the value (e.g., "Above acceptable threshold")
- evidence_quality: string — COMPLETE THIS BEFORE assigning signal_confidence. Note any of the following: the signal is inferred from document tone or pattern rather than explicitly stated; contradictory data appears elsewhere in the document; the evidence is a single ambiguous entry requiring interpretation. Write "None" only if the signal is unambiguously stated and supportable by direct reference.
- signal_confidence: string — derived from evidence_quality:
    - Hypothesis: if evidence_quality contains anything other than "None"
    - Medium: if evidence_quality is "None" AND the signal is qualitative with no specific number or count
    - High: if evidence_quality is "None" AND a specific number or count is explicitly stated with a clear label
- source: string — always "Document"
- economic_relevance: string — brief note on economic impact, or empty string
- notes: string — direct reference or quote from the document that supports this signal
{_VALENCE_FIELD}
- library_signal_id: string — ONLY include when this signal matches an entry from the SIGNAL LIBRARY block. Use the exact signal_id (e.g. "SL-19"). Omit for freely-extracted signals.

""" + _LIBRARY_INSTRUCTION + _OUTPUT_FORMAT_OBJECT

PROMPT_MAP = {
    'interview':  None,       # uses SIGNAL_EXTRACTION_PROMPT from claude.py
    'financial':  FINANCIAL_EXTRACTION_PROMPT,
    'portfolio':  PORTFOLIO_EXTRACTION_PROMPT,
    'sow':        SOW_EXTRACTION_PROMPT,
    'status':     STATUS_EXTRACTION_PROMPT,
    'resource':   RESOURCE_EXTRACTION_PROMPT,
    'delivery':   DELIVERY_DOCUMENT_EXTRACTION_PROMPT,
    'other':      None,       # uses SIGNAL_EXTRACTION_PROMPT from claude.py
}


def _build_library_block(domains: list) -> str:
    """Build a compact signal library block for injection into the user message.

    Loads Tier 1 signals for the given domains from SignalLibrary and formats
    them as a compact checklist. Returns empty string if the library is unavailable
    or empty — callers skip injection when the block is empty.
    """
    import json as _json
    from api.db.repositories.signal_library import SignalLibraryRepository

    try:
        repo = SignalLibraryRepository()
        lines = []

        for domain in domains:
            signals = [s for s in repo.get_by_domain(domain) if s['priority_tier'] == 1]
            if not signals:
                continue
            lines.append(f"\n{domain.upper()} SIGNALS:")
            for s in signals:
                sig_type = s['signal_type'].capitalize()
                lines.append(f"[{s['signal_id']}] {s['signal_name']} ({sig_type}, Tier 1)")
                defn = (s.get('definition') or '')[:200].rstrip()
                lines.append(f"  Definition: {defn}")
                if s['signal_type'] == 'numeric' and s.get('threshold_bands'):
                    try:
                        bands = _json.loads(s['threshold_bands'])
                        parts = []
                        for b in bands:
                            lo, hi, lbl = b.get('min'), b.get('max'), b.get('label', '')
                            if lo is None and hi is not None:
                                parts.append(f"{lbl} (<{hi})")
                            elif lo is not None and hi is None:
                                parts.append(f"{lbl} (>{lo})")
                            else:
                                parts.append(f"{lbl} ({lo}–{hi})")
                        lines.append(f"  Bands: {' | '.join(parts)}")
                    except (ValueError, TypeError, KeyError):
                        pass
                elif s['signal_type'] == 'maturity':
                    if s.get('maturity_levels'):
                        try:
                            levels = _json.loads(s['maturity_levels'])
                            lbls = [lv.get('label', '') for lv in levels]
                            lines.append(f"  Levels: {' | '.join(lbls)}")
                        except (ValueError, TypeError):
                            pass
                    if s.get('none_indicators'):
                        ni = s['none_indicators'][:200].rstrip()
                        lines.append(f"  None indicators: {ni}")
                lines.append('')  # blank line between signals

        if not lines:
            return ''

        header = (
            'SIGNAL LIBRARY — Check each signal below against this source.\n'
            'For matched signals add library_signal_id to found. '
            'For checked-but-not-found signals add signal_id to not_observed.\n'
        )
        return header + '\n'.join(lines)

    except Exception as e:
        logger.warning(f'_build_library_block failed: {e} — proceeding without library injection')
        return ''


def get_file_type(file_name: str) -> str:
    """Determine file type from filename for prompt routing.

    Supports two conventions:
    - New:    Interview_{role}.txt  → 'interview'
              Doc_{type}_{desc}.txt → mapped by stem substring
    - Legacy: {engagement_id}_{type}_{desc}.txt → parts[1] lookup

    Returns one of VALID_FILE_TYPES, defaulting to 'other'.
    """
    parts = Path(file_name).stem.split('_')
    prefix = parts[0].lower()

    # New convention — Interview_ prefix
    if prefix == 'interview':
        return 'interview'

    # New convention — Doc_ prefix; map remainder to file type
    if prefix == 'doc' and len(parts) >= 2:
        remainder = '_'.join(parts[1:]).lower()
        if 'financial' in remainder:
            return 'financial'
        if 'portfolio' in remainder:
            return 'portfolio'
        if 'sow' in remainder:
            return 'sow'
        if 'status' in remainder:
            return 'status'
        if 'resource' in remainder:
            return 'resource'
        if 'delivery' in remainder:
            return 'delivery'
        return 'other'

    # Legacy convention — {engagement_id}_{type}_{desc}
    if len(parts) >= 2 and parts[1].lower() in VALID_FILE_TYPES:
        return parts[1].lower()

    return 'other'


def hash_file(file_path: str) -> str:
    """Generate MD5 hash of file contents for duplicate detection."""
    with open(file_path, 'rb') as f:
        return hashlib.md5(f.read()).hexdigest()


def extract_text_from_file(file_path: str, file_name: str) -> str:
    """Extract plain UTF-8 text from a file.

    Supported formats: .txt, .docx, .xlsx, .pdf, .pptx
    Raises ValueError if the format is unsupported, the file yields no
    extractable text, or a required library is not installed.
    Never returns empty string — empty text passed to Claude produces
    hallucinated signals."""
    suffix = Path(file_name).suffix.lower()

    if suffix == '.txt':
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        if not content.strip():
            raise ValueError(f"{file_name}: file is empty")
        return content

    if suffix == '.docx':
        try:
            from docx import Document
            from docx.oxml.ns import qn
            from docx.text.paragraph import Paragraph
            from docx.table import Table as DocxTable
        except ImportError:
            logger.warning("python-docx not installed — cannot process .docx files")
            raise ValueError(f"{file_name}: python-docx not installed")

        doc = Document(file_path)
        parts = []
        # Iterate body children in document order so paragraphs and tables
        # are interleaved correctly (doc.paragraphs and doc.tables are separate
        # lists that lose interleaving order).
        for child in doc.element.body.iterchildren():
            if child.tag == qn('w:p'):
                text = Paragraph(child, doc).text.strip()
                if text:
                    parts.append(text)
            elif child.tag == qn('w:tbl'):
                for row in DocxTable(child, doc).rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    row_text = '\t'.join(c for c in cells if c)
                    if row_text:
                        parts.append(row_text)
        if not parts:
            raise ValueError(f"{file_name}: no extractable text found in Word document")
        return '\n'.join(parts)

    if suffix == '.xlsx':
        try:
            import openpyxl
        except ImportError:
            logger.warning("openpyxl not installed — cannot process .xlsx files")
            raise ValueError(f"{file_name}: openpyxl not installed")

        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"[Sheet: {sheet_name}]")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else '' for c in row]
                if any(c.strip() for c in cells):
                    parts.append('\t'.join(cells))
        wb.close()
        # Check that something beyond sheet headers was extracted
        if not any(p.strip() and not p.startswith('[Sheet:') for p in parts):
            raise ValueError(f"{file_name}: no extractable data found in Excel workbook")
        return '\n'.join(parts)

    if suffix == '.pdf':
        try:
            import pdfplumber
        except ImportError:
            logger.warning("pdfplumber not installed — cannot process .pdf files")
            raise ValueError(f"{file_name}: pdfplumber not installed")

        parts = []
        with pdfplumber.open(file_path) as pdf:
            for i, page in enumerate(pdf.pages, start=1):
                text = page.extract_text()
                if not text or not text.strip():
                    logger.debug(
                        f"{file_name}: page {i} yielded no text "
                        "(possibly scanned image) — skipped"
                    )
                    continue
                parts.append(text.strip())
        if not parts:
            raise ValueError(
                f"{file_name}: no extractable text found — PDF may be scanned or "
                "image-only. Convert to a text-searchable PDF or provide a .txt transcript."
            )
        return '\n\n'.join(parts)

    if suffix == '.pptx':
        try:
            from pptx import Presentation
        except ImportError:
            logger.warning("python-pptx not installed — cannot process .pptx files")
            raise ValueError(f"{file_name}: python-pptx not installed")

        prs = Presentation(file_path)
        parts = []
        for i, slide in enumerate(prs.slides, start=1):
            slide_parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_parts.append(text)
            if slide.has_notes_slide:
                try:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        slide_parts.append(f"[Notes: {notes_text}]")
                except (IndexError, AttributeError):
                    pass
            if slide_parts:
                parts.append(f"[Slide {i}]")
                parts.extend(slide_parts)
        if not parts:
            raise ValueError(f"{file_name}: no extractable text found in PowerPoint file")
        return '\n'.join(parts)

    raise ValueError(
        f"{file_name}: unsupported file format '{suffix}'. "
        f"Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
    )


def strip_json_fences(text: str) -> str:
    """Remove markdown code fences from Claude response.
    Handles ```json, ```, and no-fence cases.
    The CRITICAL OUTPUT FORMAT instruction reduces fence frequency but
    this remains as a safety net."""
    text = text.strip()
    if text.startswith('```json'):
        text = text[7:]
    elif text.startswith('```'):
        text = text[3:]
    if text.endswith('```'):
        text = text[:-3]
    return text.strip()


def scan_folder(folder_path: str, engagement_id: str) -> list[dict]:
    """Scan a folder for unprocessed .txt files.
    Each engagement uses a dedicated folder, so no filename prefix filter is applied.
    Returns list of dicts with path, name, type, hash."""
    if not folder_path or not os.path.exists(folder_path):
        return []

    pf_repo = ProcessedFilesRepository()
    candidates = []

    for file_name in sorted(os.listdir(folder_path)):
        if Path(file_name).suffix.lower() not in SUPPORTED_EXTENSIONS:
            continue

        file_path = os.path.join(folder_path, file_name)
        file_hash = hash_file(file_path)

        if pf_repo.already_processed(engagement_id, file_hash):
            logger.info(f"Skipping already processed file: {file_name}")
            continue

        candidates.append({
            'path': file_path,
            'name': file_name,
            'type': get_file_type(file_name),
            'hash': file_hash,
        })

    return candidates


def _parse_extraction_response(raw: str, file_name: str) -> tuple[list, list]:
    """Parse a raw signal-extraction response into (candidates, not_observed).

    Raises json.JSONDecodeError on malformed JSON or ValueError on an unexpected JSON
    shape — it NEVER silently returns empty. The caller retries once and then surfaces
    the failure, so a file's signals can never be silently dropped. (The bug this
    replaces caught the parse error, stored zero candidates, and marked the file
    processed — permanently losing every signal in that file.)"""
    clean = strip_json_fences(raw)
    parsed = json.loads(clean)  # raises json.JSONDecodeError on malformed JSON
    if isinstance(parsed, dict):
        # New format: {"found": [...], "not_observed": [...]}
        return parsed.get('found', []), parsed.get('not_observed', [])
    if isinstance(parsed, list):
        # Legacy array format (pre-Session 2 responses)
        return parsed, []
    raise ValueError(
        f"Unexpected JSON shape for {file_name}: expected object or array, "
        f"got {type(parsed).__name__}"
    )


async def _extract_signals_with_retry(file_type: str, content: str, prompt: str | None,
                                      library_block: str, file_name: str,
                                      attempts: int = 2) -> tuple[list, list]:
    """Extract signals via Claude and parse the JSON, retrying once on a parse failure.

    Malformed LLM JSON is a transient fault — calls run at temperature 1.0, so a re-call
    yields a different, usually valid, sample. Regeneration is used instead of repairing
    the bad output: repair could silently produce plausible-but-wrong JSON, trading
    silent loss for silent corruption. If every attempt fails to parse, raises — the
    caller (process_engagement_files) then records the file as failed and leaves it
    unprocessed for re-attempt, rather than losing its signals silently."""
    from api.services.claude import (
        extract_signals_from_transcript,
        extract_signals_from_document,
    )
    last_err = None
    for attempt in range(1, attempts + 1):
        if prompt is None:
            raw = await extract_signals_from_transcript(content, library_block)
        else:
            raw = await extract_signals_from_document(content, prompt, library_block)
        try:
            candidates, not_observed = _parse_extraction_response(raw, file_name)
            if attempt > 1:
                logger.info(
                    f"{file_name}: signal extraction recovered on attempt {attempt} "
                    f"after invalid JSON"
                )
            return candidates, not_observed
        except (json.JSONDecodeError, ValueError) as e:
            last_err = e
            logger.warning(
                f"{file_name}: invalid JSON from signal extraction on "
                f"attempt {attempt}/{attempts} — {e}"
            )
    raise ValueError(
        f"Signal extraction for {file_name} returned unparseable JSON after "
        f"{attempts} attempts: {last_err}"
    )


async def process_file(file_info: dict, engagement_id: str,
                       candidates_folder: str) -> dict:
    """Process a single file — extract signals via Claude and write
    candidate JSON to the candidates folder.
    Returns summary dict with file name, candidate count, and candidate file path."""
    file_type = file_info['type']
    file_path = file_info['path']
    file_name = file_info['name']

    content = extract_text_from_file(file_path, file_name)

    logger.info(f"Processing {file_type} file: {file_name} ({len(content)} chars)")

    # Build domain-filtered Tier 1 library block for this file type
    domains = DOMAIN_FILTER_MAP.get(file_type, sorted(VALID_DOMAINS))
    library_block = _build_library_block(domains)

    prompt = PROMPT_MAP.get(file_type)

    candidates, not_observed = await _extract_signals_with_retry(
        file_type, content, prompt, library_block, file_name
    )

    # Validate not_observed: only keep well-formed SL-XX IDs
    not_observed = [sid for sid in not_observed if isinstance(sid, str) and sid.startswith('SL-')]

    cleaned = []
    for item in candidates:
        if item.get('domain') not in VALID_DOMAINS:
            logger.warning(f"Invalid domain '{item.get('domain')}' in {file_name} — defaulting to {DEFAULT_DOMAIN}")
            item['domain'] = DEFAULT_DOMAIN
        if item.get('signal_confidence') not in VALID_CONFIDENCES:
            logger.warning(f"Invalid confidence '{item.get('signal_confidence')}' in {file_name} — defaulting to Medium")
            item['signal_confidence'] = 'Medium'
        # Normalize casing first — LLMs often emit 'strength'/'STRENGTH'. Title-case maps
        # those to the canonical value; only genuinely unmappable values are coerced + logged.
        v = (item.get('valence') or '').strip().capitalize()
        if v not in VALID_SIGNAL_VALENCES:
            if item.get('valence'):
                logger.warning(f"Invalid valence '{item.get('valence')}' in {file_name} — defaulting to Neutral")
            v = 'Neutral'
        item['valence'] = v
        item['interview_id'] = None
        item['source_file'] = file_name
        cleaned.append(item)

    stem = Path(file_name).stem
    os.makedirs(candidates_folder, exist_ok=True)
    candidate_file = os.path.join(candidates_folder, f"{stem}_candidates.json")

    with open(candidate_file, 'w', encoding='utf-8') as f:
        json.dump({
            'engagement_id':  engagement_id,
            'source_file':    file_name,
            'file_type':      file_type,
            'processed_date': date.today().isoformat(),
            'candidates':     cleaned,
            'not_observed':   not_observed,
        }, f, indent=2)

    logger.info(f"Wrote {len(cleaned)} candidates and {len(not_observed)} not_observed to {candidate_file}")

    return {
        'file_name':       file_name,
        'file_hash':       file_info['hash'],
        'file_type':       file_type,
        'candidate_count': len(cleaned),
        'candidate_file':  candidate_file,
        'candidates':      cleaned,
        'not_observed':    not_observed,
    }


MEDIUM_PER_DOMAIN = 3          # main-set budget for Medium-confidence signals, per domain
MAIN_CANDIDATE_CEILING = 55    # hard upper bound on the main reviewable set (volume guard)


def _select_candidates(candidates: list[dict],
                       medium_per_domain: int = MEDIUM_PER_DOMAIN,
                       ceiling: int = MAIN_CANDIDATE_CEILING,
                       strength_reserve: int = 2) -> tuple[list[dict], list[dict], int]:
    """Quality-gated selection of the main candidate set. Replaces the old fixed per-domain
    cap, which dropped material High signals on rich engagements yet admitted marginal ones
    on thin ones. Culls by QUALITY, not a fixed count:
      - Hypothesis -> hidden toggle (returned separately).
      - High: ALL kept per domain (the material core; post-dedup usually corroborated).
      - Medium: top `medium_per_domain` per domain, ranked by corroboration then evidence
        depth, reserving up to `strength_reserve` Medium slots for Strength signals (value case).
      - Global ceiling: if the main set still exceeds `ceiling` (only on very large
        engagements), trim worst-first (Medium before High, then least-corroborated) down to
        the ceiling, never below one signal per examined domain (coverage floor).
    Self-scaling: thin engagements keep few, rich ones keep their genuine material signals,
    total bounded. Returns (main, hypothesis, removed_count)."""
    rank = {'High': 2, 'Medium': 1, 'Hypothesis': 0}
    def corr(c): return c.get('_corroboration', 1)
    def rank_of(c): return (rank.get(c.get('signal_confidence', ''), 0), corr(c), len(c.get('notes') or ''))

    hypothesis = [c for c in candidates if c.get('signal_confidence') == 'Hypothesis']
    high_med = [c for c in candidates if c.get('signal_confidence') != 'Hypothesis']

    by_domain: dict[str, list[dict]] = {}
    for c in high_med:
        by_domain.setdefault(c.get('domain', ''), []).append(c)

    main: list[dict] = []
    for items in by_domain.values():
        highs = [c for c in items if c.get('signal_confidence') == 'High']
        meds = sorted((c for c in items if c.get('signal_confidence') == 'Medium'),
                      key=lambda c: (corr(c), len(c.get('notes') or '')), reverse=True)
        kept_med = meds[:medium_per_domain]
        kept_str = sum(1 for c in kept_med if c.get('valence') == 'Strength')
        if kept_str < strength_reserve:
            overflow_str = [c for c in meds[medium_per_domain:] if c.get('valence') == 'Strength']
            slots = min(strength_reserve - kept_str, len(overflow_str))
            if slots:
                ns_pos = [i for i, c in enumerate(kept_med) if c.get('valence') != 'Strength']
                drop = set(ns_pos[-slots:])
                kept_med = [c for i, c in enumerate(kept_med) if i not in drop] + overflow_str[:slots]
        main.extend(highs + kept_med)

    # global ceiling backstop — trim worst-first, but protect the best signal per domain
    if len(main) > ceiling:
        best_per_domain: dict[str, dict] = {}
        for c in main:
            d = c.get('domain', '')
            if d not in best_per_domain or rank_of(c) > rank_of(best_per_domain[d]):
                best_per_domain[d] = c
        protected = {id(c) for c in best_per_domain.values()}
        droppable = sorted((c for c in main if id(c) not in protected), key=rank_of)
        drop_ids = {id(c) for c in droppable[:len(main) - ceiling]}
        main = [c for c in main if id(c) not in drop_ids]

    removed = len(high_med) - len(main)
    for c in main + hypothesis:          # strip the transient ranking field
        c.pop('_corroboration', None)
    return main, hypothesis, removed


def _deduplicate_candidates(candidates: list[dict]) -> tuple[list[dict], int]:
    """Deduplicate candidates by (domain, normalized signal_name).
    When two candidates share the same domain and signal name (case-insensitive,
    stripped), keep the higher-confidence one and upgrade its confidence by one
    level — corroboration from a second independent file is evidence of higher
    confidence: Medium→High, Hypothesis→Medium. High stays High.
    Returns (deduped_list, count_removed)."""
    CONFIDENCE_RANK = {'High': 2, 'Medium': 1, 'Hypothesis': 0}
    CONFIDENCE_UPGRADE = {'Medium': 'High', 'Hypothesis': 'Medium', 'High': 'High'}
    seen: dict[tuple, dict] = {}
    corroborated: set[tuple] = set()
    for candidate in candidates:
        key = (
            candidate.get('domain', ''),
            candidate.get('signal_name', '').lower().strip(),
        )
        if key not in seen:
            seen[key] = candidate
        else:
            existing_rank = CONFIDENCE_RANK.get(seen[key].get('signal_confidence', ''), 0)
            incoming_rank = CONFIDENCE_RANK.get(candidate.get('signal_confidence', ''), 0)
            if incoming_rank > existing_rank:
                seen[key] = candidate
            corroborated.add(key)
    # Upgrade confidence for candidates confirmed by multiple files
    for key in corroborated:
        c = seen[key]
        old_conf = c.get('signal_confidence', 'Medium')
        new_conf = CONFIDENCE_UPGRADE.get(old_conf, old_conf)
        if new_conf != old_conf:
            logger.debug(
                f"Corroboration upgrade: {c.get('signal_name')} "
                f"({c.get('domain')}) {old_conf} → {new_conf}"
            )
            seen[key] = {**c, 'signal_confidence': new_conf}
    deduped = list(seen.values())
    return deduped, len(candidates) - len(deduped)


def _merge_cluster(members: list[dict]) -> dict:
    """Merge candidates that are the same underlying signal into one.
    Representative = highest base confidence, then longest notes (most evidence).
    Confidence = strongest member, upgraded one level if corroborated by >=2 distinct
    source files (computed ONCE here from the raw members — corroboration is not applied
    earlier). Source breadth is preserved in notes (no schema change)."""
    sources = sorted({m.get('source_file', '') for m in members if m.get('source_file')})
    corro = max(len(sources), 1)   # distinct-source count, for downstream quality ranking
    if len(members) == 1:
        return {**members[0], '_corroboration': corro}
    rank = {'High': 2, 'Medium': 1, 'Hypothesis': 0}
    upgrade = {'Hypothesis': 'Medium', 'Medium': 'High', 'High': 'High'}
    rep = max(members, key=lambda c: (rank.get(c.get('signal_confidence', ''), 0),
                                      len(c.get('notes') or '')))
    conf = max((m.get('signal_confidence', 'Hypothesis') for m in members),
               key=lambda x: rank.get(x, 0))
    merged = {**rep, '_corroboration': corro}
    if len(sources) >= 2:
        merged['signal_confidence'] = upgrade.get(conf, conf)
        note = (merged.get('notes') or '').rstrip()
        merged['notes'] = (note + f"\n[Corroborated across {len(sources)} sources: "
                           f"{', '.join(sources)}]").strip()
    else:
        merged['signal_confidence'] = conf
    return merged


async def _consolidate_candidates(all_candidates: list[dict]) -> tuple[list[dict], int, bool]:
    """Collapse duplicate candidate signals: cheap exact-name pre-grouping, then semantic
    clustering (Claude, within-domain) of the group representatives, then one merge that
    computes corroboration once from the raw members. Within-domain is enforced in CODE
    (any cross-domain cluster is split) — not left to the prompt.
    Returns (consolidated, removed_count, semantic_degraded). On semantic-clustering
    failure, falls back to exact-name dedup only (every candidate preserved — no loss),
    sets semantic_degraded=True, and logs loudly (never silently degrades)."""
    from api.services.claude import cluster_duplicate_signals
    if not all_candidates:
        return [], 0, False

    # 1. exact pre-group by (domain, normalized name) -> raw candidate indices
    groups: dict[tuple, list[int]] = {}
    order: list[tuple] = []
    for i, c in enumerate(all_candidates):
        key = (c.get('domain', ''), c.get('signal_name', '').lower().strip())
        if key not in groups:
            groups[key] = []
            order.append(key)
        groups[key].append(i)
    group_members = [groups[k] for k in order]
    reps = [all_candidates[m[0]] for m in group_members]

    # 2. semantic cluster the representatives; fall back to exact-only on any failure
    semantic_degraded = False
    try:
        clusters_of_groups = await cluster_duplicate_signals(reps)
    except Exception as exc:
        logger.warning(
            f"_consolidate_candidates: semantic dedup failed ({exc}) — falling back to "
            f"exact-name dedup only. All candidates preserved; consolidation reduced."
        )
        clusters_of_groups = [[g] for g in range(len(reps))]
        semantic_degraded = True

    # 3. enforce within-domain — split any cluster whose groups span more than one domain
    final_clusters: list[list[int]] = []
    for cluster in clusters_of_groups:
        by_dom: dict[str, list[int]] = {}
        for gid in cluster:
            by_dom.setdefault(reps[gid].get('domain', ''), []).extend(group_members[gid])
        final_clusters.extend(by_dom.values())

    # 4. one merge per final cluster (corroboration computed here, from raw members)
    consolidated = [_merge_cluster([all_candidates[i] for i in idxs]) for idxs in final_clusters]
    return consolidated, len(all_candidates) - len(consolidated), semantic_degraded


async def process_engagement_files(engagement_id: str,
                                   interviews_folder: str,
                                   documents_folder: str,
                                   candidates_folder: str) -> dict:
    """Main entry point — scan both folders, process all unprocessed files.
    Returns summary of what was processed including merged candidate file path
    and cull counts (dedup_count, hypothesis_count)."""
    all_files = []
    all_files.extend(scan_folder(interviews_folder, engagement_id))
    all_files.extend(scan_folder(documents_folder, engagement_id))

    if not all_files:
        return {'files_processed': 0, 'total_candidates': 0, 'files': [],
                'merged_candidate_file': None, 'dedup_count': 0, 'hypothesis_count': 0}

    pf_repo = ProcessedFilesRepository()
    results = []

    for file_info in all_files:
        try:
            result = await process_file(file_info, engagement_id, candidates_folder)
            pf_repo.mark_processed(
                engagement_id=engagement_id,
                file_name=file_info['name'],
                file_hash=file_info['hash'],
                file_type=file_info['type'],
                signal_count=result['candidate_count'],
            )
            results.append(result)
        except Exception as e:
            logger.error(f"Failed to process {file_info['name']}: {e}")
            results.append({
                'file_name':       file_info['name'],
                'error':           str(e),
                'candidate_count': 0,
                'candidates':      [],
                'not_observed':    [],
            })

    # Collect all candidates from all files, annotating each with its source file
    all_candidates = []
    all_not_observed: set = set()
    for result in results:
        source_file = result.get('file_name', '')
        for c in result.get('candidates', []):
            all_candidates.append({**c, 'source_file': source_file})
        all_not_observed.update(result.get('not_observed', []))

    # Consolidate across files — exact-name pre-pass + semantic dedup (within-domain)
    deduped, dedup_count, dedup_degraded = await _consolidate_candidates(all_candidates)

    # Quality-gated selection: keep all High, bound Medium, hide Hypothesis, ceiling backstop
    main_candidates, hypothesis_candidates, domain_cap_count = _select_candidates(deduped)
    hypothesis_count = len(hypothesis_candidates)

    # Signals found in any file are no longer gaps — remove from not_observed
    found_library_ids = {c.get('library_signal_id') for c in all_candidates
                         if c.get('library_signal_id')}
    final_not_observed = sorted(all_not_observed - found_library_ids)

    # Write merged candidate file
    merged_file = os.path.join(candidates_folder, f"{engagement_id}_merged_candidates.json")
    with open(merged_file, 'w', encoding='utf-8') as f:
        json.dump({
            'engagement_id':         engagement_id,
            'candidates':            main_candidates,
            'hypothesis_candidates': hypothesis_candidates,
            'dedup_count':           dedup_count,
            'domain_cap_count':      domain_cap_count,
            'hypothesis_count':      hypothesis_count,
            'semantic_dedup_degraded': dedup_degraded,
            'not_observed':          final_not_observed,
        }, f, indent=2)

    strength_count = sum(1 for c in main_candidates if c.get('valence') == 'Strength')
    logger.info(
        f"Merged candidates for {engagement_id}: "
        f"{len(main_candidates)} main, {hypothesis_count} hypothesis, "
        f"{dedup_count} duplicates removed, {domain_cap_count} removed by domain cap, "
        f"{strength_count} Strength signal(s) survived"
    )

    total_candidates = sum(r.get('candidate_count', 0) for r in results)
    return {
        'files_processed':       len(results),
        'total_candidates':      total_candidates,
        'files':                 results,
        'merged_candidate_file': merged_file,
        'dedup_count':           dedup_count,
        'domain_cap_count':      domain_cap_count,
        'hypothesis_count':      hypothesis_count,
        'semantic_dedup_degraded': dedup_degraded,
    }


def archive_candidate_files(engagement_id: str, candidates_folder: str,
                            merged_file: str) -> None:
    """Move candidate JSON files to candidates_folder/processed/ after loading.
    Archives the merged file and all individual per-file candidate JSONs for
    this engagement. Logs and continues on any failure — signals are already
    loaded and archival failure must not affect the caller."""
    if not candidates_folder or not os.path.exists(candidates_folder):
        logger.warning(f"archive_candidate_files: candidates_folder not found — {candidates_folder}")
        return

    processed_dir = os.path.join(candidates_folder, 'processed')
    try:
        os.makedirs(processed_dir, exist_ok=True)
    except OSError as e:
        logger.warning(f"archive_candidate_files: could not create processed/ dir — {e}")
        return

    files_to_archive = []

    # Merged file passed explicitly from the client
    if merged_file and os.path.exists(merged_file):
        files_to_archive.append(merged_file)

    # Individual per-file candidate JSONs — candidates_folder is per-engagement,
    # so all *_candidates.json files here (excluding the merged file) belong to
    # this engagement regardless of naming convention.
    for entry in os.scandir(candidates_folder):
        if (entry.is_file()
                and entry.name.endswith('_candidates.json')
                and entry.path != merged_file):
            files_to_archive.append(entry.path)

    for file_path in files_to_archive:
        dest = os.path.join(processed_dir, os.path.basename(file_path))
        try:
            shutil.move(file_path, dest)
            logger.info(f"Archived candidate file: {os.path.basename(file_path)}")
        except OSError as e:
            logger.warning(f"archive_candidate_files: could not move {file_path} — {e}")